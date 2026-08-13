"""Document text extraction helpers."""

import csv
import datetime
from pathlib import Path

import openpyxl
import pandas as pd
from openpyxl.utils import get_column_letter
from pptx import Presentation

from src.utils.ocr import ocr_pdf


def extract_csv_text(
    csv_path: str,
    max_rows: int = 500,
    max_chars: int = 100000,
) -> str:
    """Extract CSV content as tab-separated text for LLM analysis."""
    encodings = ["utf-8", "latin-1"]
    last_error = None

    for encoding in encodings:
        try:
            with open(csv_path, newline="", encoding=encoding) as csv_file:
                sample = csv_file.read(4096)
                csv_file.seek(0)
                try:
                    dialect = csv.Sniffer().sniff(sample)
                except csv.Error:
                    dialect = csv.excel

                rows = []
                for row_index, row in enumerate(csv.reader(csv_file, dialect)):
                    if row_index >= max_rows:
                        rows.append([f"... truncated after {max_rows} rows ..."])
                        break
                    rows.append([cell.strip() for cell in row])

            content = "\n".join("\t".join(row) for row in rows)
            return content[:max_chars]
        except UnicodeDecodeError as exc:
            last_error = exc

    raise ValueError(f"Unable to decode CSV file: {last_error}")


def _format_number(value: float, number_format: str) -> str:
    """Render a numeric cell the way Excel displays it.

    Percent-formatted cells are scaled (0.15 -> "15%") so they cannot be
    mistaken for plain ratios, and grouped formats ("#,##0") use the Vietnamese
    convention ('.' thousands, ',' decimal) to match OCR'd statement text and
    ``src.utils.formatting.to_billion_vnd``.
    """
    fmt = number_format or ""

    if "%" in fmt:
        scaled = value * 100
        text = f"{scaled:,.2f}".rstrip("0").rstrip(".")
        return _to_vietnamese_number(text or "0") + "%"

    # "#,##0" / "#,##0.00" ask Excel to group thousands.
    if "#,##" in fmt or "0,000" in fmt:
        decimals = 0
        if "." in fmt:
            decimals = len(fmt.split(".")[-1].split("_")[0].rstrip("%);\\ "))
        return _to_vietnamese_number(f"{value:,.{decimals}f}")

    if isinstance(value, float) and value.is_integer():
        return str(int(value))
    return str(value)


def _to_vietnamese_number(us_formatted: str) -> str:
    """Swap US separators for Vietnamese ones ('.' thousands, ',' decimal)."""
    return us_formatted.replace(",", "\x00").replace(".", ",").replace("\x00", ".")


def _format_cell(value: object, number_format: str) -> str:
    """Render one cell value as faithful, LLM-readable text."""
    if value is None:
        return ""

    if isinstance(value, datetime.datetime):
        if (value.hour, value.minute, value.second) == (0, 0, 0):
            return value.strftime("%Y-%m-%d")
        return value.strftime("%Y-%m-%d %H:%M")
    if isinstance(value, datetime.date):
        return value.strftime("%Y-%m-%d")
    if isinstance(value, datetime.time):
        return value.strftime("%H:%M")

    if isinstance(value, bool):
        return "TRUE" if value else "FALSE"
    if isinstance(value, (int, float)):
        return _format_number(value, number_format)

    return str(value).strip()


def _sheet_grid(worksheet) -> list[list[str]]:
    """Return the sheet as rendered text cells, with merged values propagated."""
    grid = [
        [_format_cell(cell.value, cell.number_format) for cell in row]
        for row in worksheet.iter_rows()
    ]
    if not grid:
        return []

    # A merged range stores its value only in its top-left cell; repeat it so
    # merged sub-headers/labels appear on every row and column they describe.
    # Ranges spanning the sheet's full width are cosmetic banners (report title,
    # company name), so repeating them would only add noise — keep those once.
    used_width = max((len(row) for row in grid), default=0)
    for merged in worksheet.merged_cells.ranges:
        top, left = merged.min_row - 1, merged.min_col - 1
        if top >= len(grid) or left >= len(grid[top]):
            continue
        value = grid[top][left]
        if not value:
            continue
        if merged.min_col == 1 and merged.max_col >= used_width:
            continue
        for row_index in range(top, min(merged.max_row, len(grid))):
            for col_index in range(left, merged.max_col):
                if col_index < len(grid[row_index]):
                    grid[row_index][col_index] = value
    return grid


def _clean_grid(grid: list[list[str]]) -> tuple[list[list[str]], list[int]]:
    """Drop all-empty rows/columns; return the grid and its kept column indexes."""
    if not grid:
        return [], []

    width = max(len(row) for row in grid)
    padded = [row + [""] * (width - len(row)) for row in grid]
    kept_columns = [
        index for index in range(width) if any(row[index] for row in padded)
    ]
    cleaned = [
        [row[index] for index in kept_columns]
        for row in padded
        if any(row[index] for index in kept_columns)
    ]
    return cleaned, kept_columns


def _grid_to_tsv(
    grid: list[list[str]],
    kept_columns: list[int],
    max_rows: int,
) -> str:
    """Render a cleaned grid as TSV led by its Excel column letters."""
    lines = ["\t".join(get_column_letter(index + 1) for index in kept_columns)]
    for row in grid[:max_rows]:
        # Right-trim trailing empties so short rows stop emitting runs of tabs.
        trimmed = list(row)
        while trimmed and not trimmed[-1]:
            trimmed.pop()
        lines.append("\t".join(trimmed))
    if len(grid) > max_rows:
        lines.append(f"... truncated after {max_rows} rows ...")
    return "\n".join(lines)


def extract_excel_text(
    excel_path: str,
    max_rows_per_sheet: int = 500,
    max_chars: int = 100000,
) -> str:
    """Extract XLS/XLSX workbook content as LLM-readable text grouped by sheet.

    The grid is emitted as-is (no header inference), led by Excel column letters
    so the LLM can align columns and cite a specific cell. Cells are rendered the
    way Excel displays them (percent, dates, thousands grouping).

    ``.xlsx`` is read with openpyxl in normal mode — read-only mode does not
    expose ``merged_cells.ranges``, which is needed to propagate merged headers.
    That loads the workbook into memory; ``max_rows_per_sheet`` bounds the output
    size, not peak memory. ``.xls`` (legacy) falls back to pandas + xlrd, which
    cannot recover number formats.
    """
    if Path(excel_path).suffix.lower() != ".xlsx":
        return _extract_legacy_xls_text(excel_path, max_rows_per_sheet, max_chars)

    workbook = openpyxl.load_workbook(excel_path, data_only=True)
    try:
        blocks = []
        for worksheet in workbook.worksheets:
            grid, kept_columns = _clean_grid(_sheet_grid(worksheet))
            if not grid:
                blocks.append(f"--- Sheet: {worksheet.title} (trống) ---")
                continue
            body = _grid_to_tsv(grid, kept_columns, max_rows_per_sheet)
            blocks.append(
                f"--- Sheet: {worksheet.title} "
                f"({len(grid)} dòng x {len(kept_columns)} cột) ---\n{body}"
            )
    finally:
        workbook.close()

    return _join_sheet_blocks(blocks, max_chars)


def _extract_legacy_xls_text(
    excel_path: str,
    max_rows_per_sheet: int,
    max_chars: int,
) -> str:
    """Read a legacy .xls workbook (no number-format fidelity available)."""
    workbook = pd.read_excel(
        excel_path,
        sheet_name=None,
        dtype=str,
        header=None,
        engine="xlrd",
    )

    blocks = []
    for sheet_name, dataframe in workbook.items():
        raw_grid = dataframe.fillna("").astype(str).values.tolist()
        grid, kept_columns = _clean_grid(raw_grid)
        if not grid:
            blocks.append(f"--- Sheet: {sheet_name} (trống) ---")
            continue
        body = _grid_to_tsv(grid, kept_columns, max_rows_per_sheet)
        blocks.append(
            f"--- Sheet: {sheet_name} "
            f"({len(grid)} dòng x {len(kept_columns)} cột) ---\n{body}"
        )

    return _join_sheet_blocks(blocks, max_chars)


def _join_sheet_blocks(blocks: list[str], max_chars: int) -> str:
    """Join sheet blocks, stopping with a visible marker at the char budget."""
    text = "\n\n".join(blocks)
    if len(text) <= max_chars:
        return text
    return (
        text[:max_chars].rstrip()
        + f"\n... truncated: workbook exceeds {max_chars} characters ..."
    )


def _slide_text(slide) -> list[str]:
    """Collect text from every non-title shape on a slide, tables included.

    The title shape is skipped here — it already leads the slide's block
    header (see ``extract_pptx_text``), and repeating it as the first body
    line would just echo the same text twice.
    """
    title_shape = slide.shapes.title
    title_id = title_shape.shape_id if title_shape is not None else None
    parts: list[str] = []
    for shape in slide.shapes:
        if title_id is not None and shape.shape_id == title_id:
            continue
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                parts.append(text)
        elif shape.has_table:
            for row in shape.table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    parts.append("\t".join(cells))
    return parts


def _slide_title(slide) -> str:
    """The slide's title placeholder text, if it has one."""
    title_shape = slide.shapes.title
    if title_shape is not None and title_shape.has_text_frame:
        return title_shape.text_frame.text.strip()
    return ""


def extract_pptx_text(pptx_path: str, max_chars: int = 100000) -> str:
    """Extract PowerPoint slide text (and tables) as LLM-readable text.

    Every shape with a text frame contributes its text, in shape order; a
    table's rows are rendered tab-separated, the same convention
    ``extract_excel_text`` uses for its sheet grids. One block per slide, led
    by the slide's title if it has one — mirrors the "--- Page N ---" /
    "--- Sheet: ... ---" markers used elsewhere so a citation can point at a
    specific slide.
    """
    presentation = Presentation(pptx_path)
    blocks = []
    for index, slide in enumerate(presentation.slides, start=1):
        title = _slide_title(slide)
        header = f"--- Slide {index}: {title} ---" if title else f"--- Slide {index} ---"
        lines = _slide_text(slide)
        body = "\n".join(lines) if lines else "(trống)"
        blocks.append(f"{header}\n{body}")

    return _join_sheet_blocks(blocks, max_chars)


def extract_document_text(
    file_path: str,
    ocr_timeout_seconds: float | None = None,
) -> str:
    """Extract text from supported uploaded documents."""
    spreadsheet_extensions = {".csv", ".xls", ".xlsx"}

    extension = Path(file_path).suffix.lower()
    if extension == ".pdf":
        return ocr_pdf(file_path, timeout_seconds=ocr_timeout_seconds)
    if extension == ".csv":
        return extract_csv_text(file_path)
    if extension in spreadsheet_extensions - {".csv"}:
        return extract_excel_text(file_path)
    if extension == ".pptx":
        return extract_pptx_text(file_path)
    raise ValueError(f"Unsupported file extension: {extension}")
